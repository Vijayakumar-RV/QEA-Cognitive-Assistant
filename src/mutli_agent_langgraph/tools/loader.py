import pdfplumber
from docx import Document
from pptx import Presentation
import pandas as pd
from PIL import Image
import pytesseract
import json
import os
from azure.ai.vision.imageanalysis import ImageAnalysisClient
from azure.ai.vision.imageanalysis.models import VisualFeatures
from azure.core.credentials import AzureKeyCredential

def load_pdf(file):
    text = ""
    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    return text.strip(), "application/pdf"


def load_docx(file):
    doc = Document(file)
    text = "\n".join(para.text for para in doc.paragraphs)
    return text.strip(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document"


def load_pptx(file):
    presentation = Presentation(file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip(), "application/vnd.openxmlformats-officedocument.presentationml.presentation"



def load_txt_csv(file):
    name = file.name.lower()
    if name.endswith(".csv"):
        df = pd.read_csv(file)
        text = df.to_string(index=False)
    else:
        text = file.read().decode("utf-8", errors="ignore")
    return text.strip(), "text/plain"



def load_image_ocr(file):
    image = Image.open(file).convert("RGB")
    text = pytesseract.image_to_string(image)
    return text.strip(), "image/ocr"



def load_json(file):
    try:
        data = json.load(file)
        text = json.dumps(data, indent=2)
        return text.strip(), "application/json"
    except Exception as e:
        return f"Error reading JSON: {str(e)}", "application/json"